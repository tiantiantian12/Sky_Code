"""
文档解析工具集
为 Agent 提供 PDF、Word、PPT 文档内容提取能力
"""

import os
import tempfile
from typing import Optional
from langchain_core.tools import tool


@tool
def parse_pdf(file_path: str, max_pages: int = 50) -> str:
    """解析 PDF 文件并提取文本内容。

    Args:
        file_path: PDF 文件的绝对路径
        max_pages: 最大解析页数（默认 50）
    """
    try:
        if not os.path.exists(file_path):
            return f"错误: 文件不存在 - {file_path}"
        
        ext = os.path.splitext(file_path)[1].lower()
        if ext != '.pdf':
            return f"错误: 文件不是 PDF 格式 - {ext}"

        import PyPDF2
        
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            total_pages = len(reader.pages)
            pages_to_read = min(total_pages, max_pages)
            
            result_parts = [
                f"文件: {file_path}",
                f"总页数: {total_pages}",
                f"解析页数: {pages_to_read}",
                ""
            ]
            
            for i in range(pages_to_read):
                page = reader.pages[i]
                text = page.extract_text()
                if text.strip():
                    result_parts.append(f"--- 第 {i+1} 页 ---")
                    result_parts.append(text.strip())
                    result_parts.append("")
            
            return "\n".join(result_parts)
    
    except ImportError:
        return "错误: 缺少 PyPDF2 库，请运行: pip install PyPDF2"
    except Exception as e:
        return f"错误: PDF 解析失败 - {e}"


@tool
def parse_word(file_path: str, max_paragraphs: int = 200) -> str:
    """解析 Word 文档（.docx）并提取文本内容。

    Args:
        file_path: Word 文件的绝对路径
        max_paragraphs: 最大解析段落数（默认 200）
    """
    try:
        if not os.path.exists(file_path):
            return f"错误: 文件不存在 - {file_path}"
        
        ext = os.path.splitext(file_path)[1].lower()
        if ext not in ('.docx', '.doc'):
            return f"错误: 文件不是 Word 格式 - {ext}"
        
        if ext == '.doc':
            return "错误: 暂不支持 .doc 格式，请转换为 .docx 后重试"

        from docx import Document
        
        doc = Document(file_path)
        
        result_parts = [
            f"文件: {file_path}",
            f"段落数: {len(doc.paragraphs)}",
            ""
        ]
        
        # 提取段落文本
        for i, para in enumerate(doc.paragraphs[:max_paragraphs]):
            if para.text.strip():
                result_parts.append(para.text.strip())
        
        # 提取表格内容
        if doc.tables:
            result_parts.append("")
            result_parts.append(f"--- 表格 ({len(doc.tables)} 个) ---")
            for table_idx, table in enumerate(doc.tables):
                result_parts.append(f"表格 {table_idx + 1}:")
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        result_parts.append(row_text)
                result_parts.append("")
        
        return "\n".join(result_parts)
    
    except ImportError:
        return "错误: 缺少 python-docx 库，请运行: pip install python-docx"
    except Exception as e:
        return f"错误: Word 解析失败 - {e}"


@tool
def parse_ppt(file_path: str, max_slides: int = 50) -> str:
    """解析 PowerPoint 文件（.pptx）并提取文本内容。

    Args:
        file_path: PPT 文件的绝对路径
        max_slides: 最大解析幻灯片数（默认 50）
    """
    try:
        if not os.path.exists(file_path):
            return f"错误: 文件不存在 - {file_path}"
        
        ext = os.path.splitext(file_path)[1].lower()
        if ext not in ('.pptx', '.ppt'):
            return f"错误: 文件不是 PPT 格式 - {ext}"
        
        if ext == '.ppt':
            return "错误: 暂不支持 .ppt 格式，请转换为 .pptx 后重试"

        from pptx import Presentation
        
        prs = Presentation(file_path)
        total_slides = len(prs.slides)
        slides_to_read = min(total_slides, max_slides)
        
        result_parts = [
            f"文件: {file_path}",
            f"总幻灯片数: {total_slides}",
            f"解析幻灯片数: {slides_to_read}",
            ""
        ]
        
        for i, slide in enumerate(prs.slides[:slides_to_read]):
            result_parts.append(f"--- 幻灯片 {i + 1} ---")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    result_parts.append(shape.text.strip())
                
                # 提取表格内容
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            result_parts.append(row_text)
            
            result_parts.append("")
        
        return "\n".join(result_parts)
    
    except ImportError:
        return "错误: 缺少 python-pptx 库，请运行: pip install python-pptx"
    except Exception as e:
        return f"错误: PPT 解析失败 - {e}"


@tool
def parse_document(file_path: str) -> str:
    """自动识别文档类型并解析内容。支持 PDF、Word、PPT、TXT、Markdown。

    Args:
        file_path: 文档文件的绝对路径
    """
    try:
        if not os.path.exists(file_path):
            return f"错误: 文件不存在 - {file_path}"
        
        ext = os.path.splitext(file_path)[1].lower()
        
        # 根据扩展名分发到对应的解析函数
        if ext == '.pdf':
            return parse_pdf(file_path)
        elif ext in ('.docx', '.doc'):
            return parse_word(file_path)
        elif ext in ('.pptx', '.ppt'):
            return parse_ppt(file_path)
        elif ext in ('.txt', '.md', '.markdown', '.rst', '.log'):
            return _parse_text_file(file_path)
        else:
            return f"错误: 不支持的文档格式 '{ext}'。支持: PDF, Word(.docx), PPT(.pptx), TXT, Markdown"
    
    except Exception as e:
        return f"错误: 文档解析失败 - {e}"


def _parse_text_file(file_path: str, max_size: int = 100000) -> str:
    """解析纯文本文件"""
    try:
        size = os.path.getsize(file_path)
        if size > max_size:
            return f"错误: 文件过大 ({size} bytes)，超过 {max_size} 字节限制"
        
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        
        return f"文件: {file_path}\n大小: {size} 字节\n\n{content}"
    
    except Exception as e:
        return f"错误: 文本文件读取失败 - {e}"


@tool
def extract_document_metadata(file_path: str) -> str:
    """提取文档的元数据信息（标题、作者、创建时间等）。

    Args:
        file_path: 文档文件的绝对路径
    """
    import json
    from datetime import datetime
    
    try:
        if not os.path.exists(file_path):
            return f"错误: 文件不存在 - {file_path}"
        
        ext = os.path.splitext(file_path)[1].lower()
        
        # 基本文件信息
        stat = os.stat(file_path)
        metadata = {
            "文件名": os.path.basename(file_path),
            "文件大小": f"{stat.st_size} 字节",
            "修改时间": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
            "创建时间": datetime.fromtimestamp(stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S"),
        }
        
        # 根据文件类型提取特定元数据
        if ext == '.pdf':
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                if reader.metadata:
                    metadata["标题"] = reader.metadata.title or "无"
                    metadata["作者"] = reader.metadata.author or "无"
                    metadata["主题"] = reader.metadata.subject or "无"
                    metadata["创建工具"] = reader.metadata.creator or "无"
                metadata["页数"] = len(reader.pages)
        
        elif ext == '.docx':
            from docx import Document
            doc = Document(file_path)
            core = doc.core_properties
            metadata["标题"] = core.title or "无"
            metadata["作者"] = core.author or "无"
            metadata["主题"] = core.subject or "无"
            metadata["关键词"] = core.keywords or "无"
            metadata["段落数"] = len(doc.paragraphs)
            metadata["表格数"] = len(doc.tables)
        
        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            core = prs.core_properties
            metadata["标题"] = core.title or "无"
            metadata["作者"] = core.author or "无"
            metadata["主题"] = core.subject or "无"
            metadata["幻灯片数"] = len(prs.slides)
        
        return json.dumps(metadata, ensure_ascii=False, indent=2)
    
    except ImportError as e:
        return f"错误: 缺少必要的库 - {e}"
    except Exception as e:
        return f"错误: 元数据提取失败 - {e}"
